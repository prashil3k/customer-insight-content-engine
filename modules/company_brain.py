import json
import time
import requests
from bs4 import BeautifulSoup
import config
from modules.model_manager import create_message


def get_company_brief() -> dict:
    if config.COMPANY_BRIEF_PATH.exists():
        return json.loads(config.COMPANY_BRIEF_PATH.read_text())
    return {}


def brief_as_prompt_context() -> str:
    brief = get_company_brief()
    if not brief:
        return ""
    lines = [
        "=== STORYLANE COMPANY INTELLIGENCE ===",
        f"What Storylane does: {brief.get('what_we_do', '')}",
        f"ICP: {brief.get('icp', '')}",
        f"Core capabilities: {', '.join(brief.get('capabilities', []))}",
        f"Key differentiators: {', '.join(brief.get('differentiators', []))}",
        f"Use cases we own: {', '.join(brief.get('use_cases', []))}",
        f"Main competitors: {', '.join(brief.get('competitors', []))}",
        f"Tone of voice: {brief.get('tone_of_voice', '')}",
        f"What good content means for us: {brief.get('content_north_star', '')}",
        "=== END COMPANY INTELLIGENCE ===",
    ]
    return "\n".join(lines)


def extract_text_from_file(filepath: str, filename: str) -> str:
    ext = filename.rsplit(".", 1)[-1].lower()
    try:
        if ext == "pdf":
            import pdfplumber
            with pdfplumber.open(filepath) as pdf:
                pages = [p.extract_text() or "" for p in pdf.pages[:30]]
            return "\n".join(pages)[:40000]
        elif ext in ("pptx", "ppt"):
            from pptx import Presentation
            prs = Presentation(filepath)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
            return "\n".join(parts)[:40000]
        elif ext in ("xlsx", "xls"):
            import openpyxl
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            parts = []
            for sheet in wb.worksheets[:5]:
                parts.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(max_row=200, values_only=True):
                    row_str = " | ".join(str(c) for c in row if c is not None)
                    if row_str.strip():
                        parts.append(row_str)
            return "\n".join(parts)[:40000]
        elif ext == "csv":
            import csv, io
            with open(filepath, newline="", errors="replace") as f:
                reader = csv.reader(f)
                rows = [" | ".join(r) for r in reader]
            return "\n".join(rows[:500])
        else:
            return open(filepath, errors="replace").read()[:40000]
    except Exception as e:
        return f"[Could not extract {filename}: {e}]"


def _fetch_url(url: str) -> str:
    try:
        headers = {"User-Agent": "Mozilla/5.0 (compatible; StorylaneCE/1.0)"}
        r = requests.get(url, headers=headers, timeout=15)
        r.raise_for_status()
        soup = BeautifulSoup(r.text, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        lines = [l for l in text.splitlines() if len(l.strip()) > 30]
        return "\n".join(lines[:300])
    except Exception as e:
        return f"[Could not fetch {url}: {e}]"


def scan_company_intelligence(urls: list, raw_docs: list = None, progress_cb=None) -> dict:
    raw_docs = raw_docs or []

    def _progress(msg):
        if progress_cb:
            progress_cb(msg)

    _progress("Fetching URLs...")
    url_contents = []
    for url in urls:
        _progress(f"Fetching {url}")
        content = _fetch_url(url)
        url_contents.append(f"--- SOURCE: {url} ---\n{content}")
        time.sleep(0.5)

    all_content = "\n\n".join(url_contents + raw_docs)

    _progress("Distilling intelligence with Claude...")
    prompt = f"""You are building a company intelligence brief for Storylane. Below is raw content from Storylane's website and documents.

Extract and distill ONLY the essential intelligence a content marketer needs. Be specific and concrete — no vague language.

RAW CONTENT:
{all_content[:80000]}

Return a JSON object with exactly these keys:
- "what_we_do": One crisp sentence describing Storylane's core product (max 30 words)
- "icp": Who they sell to — specific role, company stage, industry, pain trigger (2-3 sentences)
- "capabilities": Array of 6-10 core product capability clusters (not feature names, e.g. "Interactive demo creation without engineering")
- "differentiators": Array of 4-6 key differentiators vs competitors (specific, not generic)
- "use_cases": Array of 6-10 primary use cases they own (e.g. "Replacing free trials with interactive demos")
- "competitors": Array of main competitors by name
- "pricing_tiers": Brief description of pricing model (1-2 sentences, omit if not clear)
- "tone_of_voice": How Storylane communicates — style, personality, what they avoid (2-3 sentences)
- "content_north_star": What good content means for Storylane — what it should achieve, what authority looks like (2-3 sentences)
- "eeat_angles": Array of EEAT signal types that make sense for Storylane (e.g. "Customer proof with specific metrics", "Founder POV on demo psychology")
- "backlinkable_asset_types": Array of content types that could earn backlinks in this space

Return ONLY the JSON object, no other text."""

    response = create_message("sonnet", max_tokens=2000, messages=[{"role": "user", "content": prompt}])

    raw = response.content[0].text.strip()
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
    raw = raw.strip().rstrip("```").strip()

    brief = json.loads(raw)
    brief["scanned_urls"] = urls
    brief["scanned_at"] = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())

    config.COMPANY_BRIEF_PATH.write_text(json.dumps(brief, indent=2))
    _progress("Company intelligence saved.")
    return brief
